from __future__ import annotations

import io
from pathlib import Path
from typing import Optional

from app.schemas.extraction_types import ExtractionTuning
from app.services.extraction_service import ExtractionService


# (image_bytes, location_label) — one entry per candidate crop found in a file.
Candidate = tuple[bytes, str]


def dispatch_file(
    filename: str,
    data: bytes,
    tuning: Optional[ExtractionTuning] = None,
) -> list[Candidate]:
    """Turn one uploaded file into a list of candidate structure crops.

    Every format funnels into the same OpenCV detector (ExtractionService): PDFs
    are rastered page-by-page, standalone images are decoded directly, and each
    image embedded in a PPTX slide is decoded and detected. Raises ValueError for
    unsupported extensions so the caller can record a per-file error row.
    """
    ext = Path(filename).suffix.lower()

    if ext == ".pdf":
        patent = ExtractionService.extract_candidate_compounds(data, patent_id="upload", tuning=tuning)
        return [(c.image_bytes, f"page {c.page}") for c in patent.get_compounds()]

    if ext in {".png", ".jpg", ".jpeg", ".heic", ".heif"}:
        # HEIC/HEIF (e.g. iPhone photos) can't be decoded by OpenCV, so normalize to
        # PNG bytes first; PNG/JPG pass through unchanged.
        image_data = _heic_to_png(data) if ext in {".heic", ".heif"} else data
        compounds = ExtractionService.detect_in_image_bytes(image_data, tuning=tuning)
        # Detection is tuned for multi-structure pages; a clean single-structure
        # image may yield nothing. Fall back to the whole image as one candidate so
        # such inputs aren't silently dropped.
        if not compounds:
            return [(image_data, "image")]
        return [(c.image_bytes, "image") for c in compounds]

    if ext == ".pptx":
        return _from_pptx(data, tuning)

    if ext == ".ppt":
        # Legacy binary PowerPoint is not readable by python-pptx.
        raise ValueError("Legacy .ppt is not supported; please export to .pptx")

    raise ValueError(f"Unsupported file type: {ext or '(none)'}")


def _heic_to_png(data: bytes) -> bytes:
    """Decode HEIC/HEIF bytes and re-encode as PNG so the rest of the pipeline
    (OpenCV detection, MolScribe) can read them."""
    try:
        import pillow_heif
        from PIL import Image
    except ImportError as exc:  # pragma: no cover - dependency guard
        raise RuntimeError(f"pillow-heif is required for HEIC files: {exc}") from exc

    pillow_heif.register_heif_opener()  # idempotent: teaches Pillow to open HEIF
    image = Image.open(io.BytesIO(data)).convert("RGB")
    out = io.BytesIO()
    image.save(out, format="PNG")
    return out.getvalue()


def _from_pptx(data: bytes, tuning: Optional[ExtractionTuning]) -> list[Candidate]:
    """Extract each embedded picture from a .pptx and run detection on it.

    python-pptx exposes slides -> shapes; picture shapes carry the original image
    bytes via ``shape.image.blob``. A slide may hold several structures, so each
    picture is passed through the detector (which can yield multiple crops).
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError as exc:  # pragma: no cover - dependency guard
        raise RuntimeError(f"python-pptx is required for .pptx files: {exc}") from exc

    presentation = Presentation(io.BytesIO(data))
    candidates: list[Candidate] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
            blob = shape.image.blob
            compounds = ExtractionService.detect_in_image_bytes(
                blob, location_index=slide_index, tuning=tuning
            )
            # Same fallback as standalone images: if detection finds nothing in a
            # picture, treat the whole picture as a single candidate.
            if not compounds:
                candidates.append((blob, f"slide {slide_index}"))
            else:
                candidates.extend((c.image_bytes, f"slide {slide_index}") for c in compounds)
    return candidates
